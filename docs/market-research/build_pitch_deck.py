#!/usr/bin/env python3
"""Generate NeuralMind corporate pitch deck PPTX."""

import sys

sys.path.insert(0, "/home/dtfrost/neuralmind/docs/market-research")

from pathlib import Path

from intel import (
    COMPETITORS,
    MARKET,
    PRICING,
    REFACTOR_PRIORITIES,
    SEGMENTS,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme Colors ──────────────────────────────────────────────────────────
NAVY = RGBColor(0x1A, 0x1F, 0x2E)  # Primary dark background
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x3A)  # Secondary dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x00, 0x66, 0xCC)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MID_GRAY = RGBColor(0x88, 0x88, 0x99)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
AMBER = RGBColor(0xFF, 0xC1, 0x07)

# ── Slide dimensions (16:9) ───────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Typography helpers ────────────────────────────────────────────────────
def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=18,
    color=WHITE,
    bold=False,
    alignment=PP_ALIGN.LEFT,
    font_name="Calibri",
):
    """Add a textbox to a slide."""
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tx_box


def add_shape(slide, left, top, width, height, fill=ACCENT_BLUE):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def style_cell(
    cell, text, font_size=11, color=WHITE, bold=False, fill=None, alignment=PP_ALIGN.LEFT
):
    """Style a table cell."""
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
        paragraph.font.bold = bold
        paragraph.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


def add_footer(slide, text="NeuralMind — Confidential"):
    """Add a footer to a slide."""
    add_textbox(
        slide,
        Inches(0.3),
        Inches(7.1),
        Inches(12.7),
        Inches(0.3),
        text,
        font_size=8,
        color=MID_GRAY,
    )


def add_slide_number(slide, num, total):
    """Add slide number."""
    add_textbox(
        slide,
        Inches(12.5),
        Inches(7.1),
        Inches(0.7),
        Inches(0.3),
        f"{num}/{total}",
        font_size=8,
        color=MID_GRAY,
        alignment=PP_ALIGN.RIGHT,
    )


# ── Main deck builder ────────────────────────────────────────────────────
def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    total_slides = 10

    # ── SLIDE 1: Title ────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    # Accent bar
    add_shape(slide, Inches(0.8), Inches(2.2), Inches(0.08), Inches(1.8), fill=ACCENT_BLUE)

    add_textbox(
        slide,
        Inches(1.2),
        Inches(2.0),
        Inches(11),
        Inches(0.8),
        "NeuralMind",
        font_size=52,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(1.2),
        Inches(2.9),
        Inches(11),
        Inches(0.6),
        "Local-First Graph Intelligence for Engineering Teams",
        font_size=24,
        color=ACCENT_BLUE,
    )
    add_textbox(
        slide,
        Inches(1.2),
        Inches(3.7),
        Inches(11),
        Inches(0.4),
        "Privacy-preserving AI that indexes your codebase — no code leaves your premise.",
        font_size=16,
        color=LIGHT_GRAY,
    )

    add_textbox(
        slide,
        Inches(1.2),
        Inches(5.5),
        Inches(5),
        Inches(0.3),
        "Corporate Pitch Deck — July 2026",
        font_size=12,
        color=MID_GRAY,
    )
    add_textbox(
        slide,
        Inches(1.2),
        Inches(5.9),
        Inches(5),
        Inches(0.3),
        "Confidential — For Investor Discussion",
        font_size=10,
        color=MID_GRAY,
    )

    add_slide_number(slide, 1, total_slides)

    # ── SLIDE 2: Problem ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "Cloud AI Creates New Risks for Engineering Orgs",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.2),
        Inches(11.5),
        Inches(0.4),
        "The shift to cloud AI coding assistants introduces critical vulnerabilities:",
        font_size=16,
        color=LIGHT_GRAY,
    )

    problems = [
        (
            "Data Sovereignty Risk",
            "Source code sent to third-party APIs — violating compliance requirements (SOC2, HIPAA, ITAR).",
        ),
        (
            "Cost Escalation",
            "Cloud AI pricing at $19–39/user/mo scales linearly — a 500-engineer org faces $1.5M+/yr in API costs.",
        ),
        (
            "Vendor Lock-in",
            "Proprietary models and formats create switching costs; code intelligence is trapped in one ecosystem.",
        ),
        (
            "Project Failure Rate",
            f"{MARKET['ai_project_cancellation_2027']:.0%} of AI projects will be cancelled by 2027 due to data governance and cost (Gartner).",
        ),
        (
            "Knowledge Loss",
            "When engineers leave, cloud AI context leaves with them — no persistent institutional memory.",
        ),
    ]

    y = Inches(1.9)
    for title, desc in problems:
        add_shape(slide, Inches(0.8), y, Inches(0.06), Inches(0.5), fill=ACCENT_BLUE)
        add_textbox(
            slide,
            Inches(1.0),
            y,
            Inches(3.5),
            Inches(0.3),
            title,
            font_size=14,
            color=ACCENT_BLUE,
            bold=True,
        )
        add_textbox(
            slide,
            Inches(1.0),
            y + Inches(0.28),
            Inches(11.0),
            Inches(0.25),
            desc,
            font_size=11,
            color=LIGHT_GRAY,
        )
        y += Inches(0.85)

    add_footer(slide)
    add_slide_number(slide, 2, total_slides)

    # ── SLIDE 3: Solution ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "NeuralMind: Graph Intelligence That Never Leaves Your Premise",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.1),
        Inches(11.5),
        Inches(0.4),
        "Local-first architecture with semantic graph indexing — AI that works for you, not your cloud provider.",
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Two columns
    # Left: Engineering value
    add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(0.4), fill=DARK_GRAY)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(1.82),
        Inches(5.4),
        Inches(0.35),
        "Engineering Value",
        font_size=16,
        color=ACCENT_BLUE,
        bold=True,
    )

    eng_items = [
        "On-device graph indexing — zero code egress",
        "Semantic search across 8000+ node codebases",
        "Self-documenting via synapse learning",
        "Multi-language: Python, TS, Go, Rust, Java, C/C++",
        "IDE-agnostic via MCP protocol",
    ]
    y = Inches(2.3)
    for item in eng_items:
        add_textbox(
            slide,
            Inches(1.0),
            y,
            Inches(5.4),
            Inches(0.25),
            f"▸  {item}",
            font_size=12,
            color=WHITE,
        )
        y += Inches(0.42)

    # Right: Business value
    add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(0.4), fill=DARK_GRAY)
    add_textbox(
        slide,
        Inches(7.0),
        Inches(1.82),
        Inches(5.4),
        Inches(0.35),
        "Business Value",
        font_size=16,
        color=GREEN,
        bold=True,
    )

    biz_items = [
        "Engineering productivity metrics via audit trail",
        "Code knowledge persists through turnover",
        "Reduced onboarding time for new engineers",
        "Compliance: full data sovereignty",
        "Lower TCO vs cloud AI (no per-token API costs)",
    ]
    y = Inches(2.3)
    for item in biz_items:
        add_textbox(
            slide,
            Inches(7.0),
            y,
            Inches(5.4),
            Inches(0.25),
            f"▸  {item}",
            font_size=12,
            color=WHITE,
        )
        y += Inches(0.42)

    add_footer(slide)
    add_slide_number(slide, 3, total_slides)

    # ── SLIDE 4: Market ───────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "A $7.65B Market Growing to $22.2B by 2030",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    # Big numbers row
    box_data = [
        ("$7.65B", "2025 Market Size", ACCENT_BLUE),
        ("$22.2B", "2030 Projected", GREEN),
        ("23.8%", "CAGR", AMBER),
    ]
    x = Inches(0.8)
    for num, label, color in box_data:
        add_shape(slide, x, Inches(1.3), Inches(3.5), Inches(1.5), fill=DARK_GRAY)
        add_shape(slide, x, Inches(1.3), Inches(3.5), Inches(0.06), fill=color)
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(1.45),
            Inches(3.1),
            Inches(0.6),
            num,
            font_size=36,
            color=color,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(2.1),
            Inches(3.1),
            Inches(0.3),
            label,
            font_size=12,
            color=LIGHT_GRAY,
        )
        x += Inches(3.9)

    # Market details
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.1),
        Inches(11.5),
        Inches(0.4),
        "Key Market Dynamics",
        font_size=18,
        color=WHITE,
        bold=True,
    )

    dynamics = [
        f"Cloud dominates at {MARKET['cloud_share_2025']:.1%} of market — but on-prem growing faster at {MARKET['onprem_cagr']:.1%} CAGR",
        f"Enterprise AI market already ${MARKET['enterprise_ai_2024']}B in 2024 — code generation is {MARKET['code_generation_pct']:.0%} of use cases",
        f"{MARKET['developer_adoption_2025']:.0%} of developers already use AI coding tools — switching costs are real",
        f"Gartner: {MARKET['gartner_2028_developer_adoption']:.0%} of developers will use AI coding assistants by 2028",
        "Privacy regulations (GDPR, state laws) accelerating on-prem demand",
    ]
    y = Inches(3.6)
    for d in dynamics:
        add_textbox(
            slide,
            Inches(1.0),
            y,
            Inches(11.3),
            Inches(0.28),
            f"•  {d}",
            font_size=12,
            color=LIGHT_GRAY,
        )
        y += Inches(0.42)

    add_footer(slide)
    add_slide_number(slide, 4, total_slides)

    # ── SLIDE 5: How It Works ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "Graph Intelligence: How NeuralMind Works",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.1),
        Inches(11.5),
        Inches(0.4),
        "A semantic graph of your codebase — built locally, queried intelligently, always up to date.",
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Pipeline steps
    steps = [
        (
            "1. Parse",
            "Multi-language graph parsing via tree-sitter/SCIP\n8000+ node capacity, incremental updates",
        ),
        ("2. Index", "Pluggable vector backends\nChromaDB / TurboVec / custom"),
        ("3. Query", "Semantic search + graph traversal\nMCP protocol, IDE-agnostic"),
        ("4. Learn", "Self-improving synapse layer\nAudit trail, knowledge persistence"),
    ]

    x = Inches(0.6)
    for title, desc in steps:
        # Box
        add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(3.5), fill=DARK_GRAY)
        add_shape(slide, x, Inches(1.8), Inches(2.8), Inches(0.06), fill=ACCENT_BLUE)
        # Title
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(2.0),
            Inches(2.4),
            Inches(0.4),
            title,
            font_size=20,
            color=ACCENT_BLUE,
            bold=True,
        )
        # Description
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(2.5),
            Inches(2.4),
            Inches(2.5),
            desc,
            font_size=12,
            color=LIGHT_GRAY,
        )
        x += Inches(3.1)

    # Bottom: Key differentiators
    add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.05), fill=DARK_GRAY)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.8),
        Inches(11.5),
        Inches(0.3),
        "Key Differentiators",
        font_size=14,
        color=WHITE,
        bold=True,
    )

    diffs = [
        "Zero code egress — all processing on-prem",
        "Open architecture — no vendor lock-in",
        "Self-learning — improves with usage",
        "Enterprise-ready — SSO, RBAC, audit trail",
    ]
    x = Inches(0.8)
    for d in diffs:
        add_textbox(
            slide, x, Inches(6.2), Inches(2.7), Inches(0.3), f"✓ {d}", font_size=11, color=GREEN
        )
        x += Inches(2.9)

    add_footer(slide)
    add_slide_number(slide, 5, total_slides)

    # ── SLIDE 6: Competitive Landscape ────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "Competitive Landscape",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    # Competitor table
    rows = len(COMPETITORS) + 1
    cols = 5
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.1), Inches(12.3), Inches(4.5)
    )
    tbl = tbl_shape.table

    # Column widths
    tbl.columns[0].width = Inches(1.8)  # Name
    tbl.columns[1].width = Inches(2.2)  # Type
    tbl.columns[2].width = Inches(2.8)  # Strength
    tbl.columns[3].width = Inches(2.8)  # Weakness
    tbl.columns[4].width = Inches(2.5)  # Pricing

    headers = ["Player", "Type", "Strength", "Weakness", "Pricing"]
    for c, h in enumerate(headers):
        style_cell(
            tbl.cell(0, c),
            h,
            font_size=11,
            color=WHITE,
            bold=True,
            fill=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER,
        )

    for r, comp in enumerate(COMPETITORS, start=1):
        is_us = comp["name"] == "NeuralMind"
        row_fill = DARK_GRAY if is_us else None
        name_color = ACCENT_BLUE if is_us else WHITE
        style_cell(
            tbl.cell(r, 0), comp["name"], font_size=10, color=name_color, bold=is_us, fill=row_fill
        )
        style_cell(tbl.cell(r, 1), comp["type"], font_size=10, color=LIGHT_GRAY, fill=row_fill)
        style_cell(tbl.cell(r, 2), comp["strength"], font_size=10, color=LIGHT_GRAY, fill=row_fill)
        style_cell(tbl.cell(r, 3), comp["weakness"], font_size=10, color=LIGHT_GRAY, fill=row_fill)
        style_cell(tbl.cell(r, 4), comp["pricing"], font_size=10, color=LIGHT_GRAY, fill=row_fill)

    # Positioning statement
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.9),
        Inches(11.5),
        Inches(0.3),
        "NeuralMind's Position",
        font_size=14,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.25),
        Inches(11.5),
        Inches(0.4),
        "Only local-first, graph-intelligence platform with zero code egress — combining the privacy of on-prem with the intelligence of cloud AI.",
        font_size=12,
        color=ACCENT_BLUE,
    )

    add_footer(slide)
    add_slide_number(slide, 6, total_slides)

    # ── SLIDE 7: Refactoring Priorities ───────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "Refactoring Priorities — Engineering Roadmap",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.1),
        Inches(11.5),
        Inches(0.4),
        "Codebase hardening to support enterprise scale and security requirements.",
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Priority table
    rows = len(REFACTOR_PRIORITIES) + 1
    cols = 4
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.7), Inches(12.3), Inches(3.5)
    )
    tbl = tbl_shape.table

    tbl.columns[0].width = Inches(4.0)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(2.0)
    tbl.columns[3].width = Inches(4.3)

    headers = ["Item", "Impact", "Effort", "Rationale"]
    for c, h in enumerate(headers):
        style_cell(
            tbl.cell(0, c),
            h,
            font_size=11,
            color=WHITE,
            bold=True,
            fill=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER,
        )

    impact_colors = {"HIGH": AMBER, "MED": ACCENT_BLUE, "LOW": MID_GRAY}
    effort_colors = {"HIGH": AMBER, "MED": ACCENT_BLUE, "LOW": GREEN}

    for r, item in enumerate(REFACTOR_PRIORITIES, start=1):
        style_cell(tbl.cell(r, 0), item["item"], font_size=11, color=WHITE)
        style_cell(
            tbl.cell(r, 1),
            item["impact"],
            font_size=11,
            color=impact_colors.get(item["impact"], WHITE),
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        style_cell(
            tbl.cell(r, 2),
            item["effort"],
            font_size=11,
            color=effort_colors.get(item["effort"], WHITE),
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )
        style_cell(tbl.cell(r, 3), item["rationale"], font_size=11, color=LIGHT_GRAY)

    # Bottom note
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.5),
        Inches(11.5),
        Inches(0.3),
        "Immediate Actions",
        font_size=14,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.85),
        Inches(11.5),
        Inches(0.4),
        "Remove regex-experiment.py (dangerous source mutation) → Decompose core.py (1811 lines) → Modularize ir.py (963 lines)",
        font_size=12,
        color=ACCENT_BLUE,
    )

    add_footer(slide)
    add_slide_number(slide, 7, total_slides)

    # ── SLIDE 8: Go-to-Market ─────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "Go-to-Market Strategy",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.1),
        Inches(11.5),
        Inches(0.4),
        "Land via OSS community, expand through enterprise sales — bottoms-up adoption with top-down expansion.",
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Pricing tiers
    tiers = [
        ("OSS", PRICING["oss"]["price"], PRICING["oss"]["features"], MID_GRAY),
        ("Team", PRICING["team"]["price"], PRICING["team"]["features"], ACCENT_BLUE),
        ("Enterprise", PRICING["enterprise"]["price"], PRICING["enterprise"]["features"], GREEN),
    ]

    x = Inches(0.6)
    for name, price, features, color in tiers:
        add_shape(slide, x, Inches(1.8), Inches(3.7), Inches(2.8), fill=DARK_GRAY)
        add_shape(slide, x, Inches(1.8), Inches(3.7), Inches(0.06), fill=color)
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(2.0),
            Inches(3.3),
            Inches(0.35),
            name,
            font_size=18,
            color=color,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(2.4),
            Inches(3.3),
            Inches(0.35),
            price,
            font_size=22,
            color=WHITE,
            bold=True,
        )
        add_textbox(
            slide,
            x + Inches(0.2),
            Inches(2.85),
            Inches(3.3),
            Inches(1.5),
            features,
            font_size=11,
            color=LIGHT_GRAY,
        )
        x += Inches(4.1)

    # Target segments
    add_textbox(
        slide,
        Inches(0.8),
        Inches(4.9),
        Inches(11.5),
        Inches(0.3),
        "Target Segments",
        font_size=14,
        color=WHITE,
        bold=True,
    )

    seg_text = "  •  ".join([f"{s['name']}: {s['pain']}" for s in SEGMENTS])
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.25),
        Inches(11.5),
        Inches(0.8),
        seg_text,
        font_size=11,
        color=LIGHT_GRAY,
    )

    add_footer(slide)
    add_slide_number(slide, 8, total_slides)

    # ── SLIDE 9: Financial Model ──────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "5-Year Financial Model — Conservative",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    # Financial table
    rows = 7  # header + 5 years + cumulative
    cols = 6
    tbl_shape = slide.shapes.add_table(
        rows, cols, Inches(0.5), Inches(1.1), Inches(12.3), Inches(4.0)
    )
    tbl = tbl_shape.table

    tbl.columns[0].width = Inches(1.5)
    for i in range(1, 6):
        tbl.columns[i].width = Inches(2.16)

    headers = ["Metric", "Y1", "Y2", "Y3", "Y4", "Y5"]
    for c, h in enumerate(headers):
        style_cell(
            tbl.cell(0, c),
            h,
            font_size=11,
            color=WHITE,
            bold=True,
            fill=ACCENT_BLUE,
            alignment=PP_ALIGN.CENTER,
        )

    # Data rows
    metrics = [
        ("Revenue", ["$0", "$45K", "$320K", "$950K", "$2.4M"]),
        ("Cost", ["$85K", "$180K", "$380K", "$620K", "$1.1M"]),
        ("Net", ["-$85K", "-$135K", "-$60K", "$330K", "$1.3M"]),
        ("Users", ["0", "50", "250", "750", "2,000"]),
        ("Paying Teams", ["0", "5", "10", "30", "75"]),
        ("Notes", ["OSS Launch", "Pilots", "Ent. Launch", "Scale", "Category Lead."]),
    ]

    for r, (metric, values) in enumerate(metrics, start=1):
        style_cell(tbl.cell(r, 0), metric, font_size=11, color=WHITE, bold=True)
        for c, val in enumerate(values, start=1):
            is_net = metric == "Net"
            color = GREEN if is_net and not val.startswith("-") else (AMBER if is_net else WHITE)
            style_cell(tbl.cell(r, c), val, font_size=11, color=color, alignment=PP_ALIGN.CENTER)

    # Key metrics summary
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.3),
        Inches(11.5),
        Inches(0.3),
        "Key Takeaways",
        font_size=14,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.65),
        Inches(11.5),
        Inches(0.5),
        "Break-even in Year 4  •  $2.4M revenue by Year 5  •  2,000 users  •  Conservative 23.8% market CAGR assumption",
        font_size=12,
        color=ACCENT_BLUE,
    )

    add_footer(slide)
    add_slide_number(slide, 9, total_slides)

    # ── SLIDE 10: Team/Ops Ask ────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(0.08), Inches(0.5), fill=ACCENT_BLUE)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(0.3),
        Inches(11),
        Inches(0.6),
        "The Ask — Team & Operations",
        font_size=32,
        color=WHITE,
        bold=True,
    )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.1),
        Inches(11.5),
        Inches(0.4),
        "Building the team and infrastructure to capture the $22.2B local-first AI opportunity.",
        font_size=14,
        color=LIGHT_GRAY,
    )

    # Two columns
    # Left: Hiring Plan
    add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.8), Inches(0.4), fill=DARK_GRAY)
    add_textbox(
        slide,
        Inches(1.0),
        Inches(1.82),
        Inches(5.4),
        Inches(0.35),
        "Hiring Plan (Year 1-2)",
        font_size=16,
        color=ACCENT_BLUE,
        bold=True,
    )

    hires = [
        ("Engineering", "2 senior backend, 1 ML/inference, 1 security"),
        ("Product", "1 product manager, 1 developer advocate"),
        ("Sales", "1 enterprise AE, 1 SDR"),
        ("Operations", "1 ops lead, 1 compliance/security"),
    ]
    y = Inches(2.3)
    for role, detail in hires:
        add_textbox(
            slide,
            Inches(1.0),
            y,
            Inches(1.2),
            Inches(0.25),
            role,
            font_size=11,
            color=ACCENT_BLUE,
            bold=True,
        )
        add_textbox(
            slide, Inches(2.3), y, Inches(4.3), Inches(0.25), detail, font_size=11, color=LIGHT_GRAY
        )
        y += Inches(0.45)

    # Right: Use of Funds
    add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(0.4), fill=DARK_GRAY)
    add_textbox(
        slide,
        Inches(7.0),
        Inches(1.82),
        Inches(5.4),
        Inches(0.35),
        "Use of Funds",
        font_size=16,
        color=GREEN,
        bold=True,
    )

    funds = [
        ("Engineering & R&D", "45% — Core platform, graph intelligence, security hardening"),
        ("Sales & Marketing", "25% — Enterprise sales, developer relations, community"),
        ("Infrastructure", "15% — On-prem deployment tooling, CI/CD, security audits"),
        ("Operations", "15% — Legal, compliance, finance, admin"),
    ]
    y = Inches(2.3)
    for category, detail in funds:
        add_textbox(
            slide,
            Inches(7.0),
            y,
            Inches(1.8),
            Inches(0.25),
            category,
            font_size=11,
            color=GREEN,
            bold=True,
        )
        add_textbox(
            slide, Inches(8.9), y, Inches(3.7), Inches(0.25), detail, font_size=11, color=LIGHT_GRAY
        )
        y += Inches(0.45)

    # Milestones
    add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.05), fill=DARK_GRAY)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(4.7),
        Inches(11.5),
        Inches(0.3),
        "12-Month Milestones",
        font_size=14,
        color=WHITE,
        bold=True,
    )

    milestones = [
        "OSS launch + 500 GitHub stars",
        "3 enterprise pilot customers",
        "SOC2 Type I certification",
        "Core.py decomposition complete",
        "$45K ARR from Team tier",
    ]
    x = Inches(0.8)
    for m in milestones:
        add_textbox(
            slide,
            x,
            Inches(5.1),
            Inches(2.2),
            Inches(0.3),
            f"◆ {m}",
            font_size=10,
            color=ACCENT_BLUE,
        )
        x += Inches(2.3)

    # Closing
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.2),
        Inches(11.5),
        Inches(0.4),
        "NeuralMind — Local-First Graph Intelligence",
        font_size=18,
        color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.6),
        Inches(11.5),
        Inches(0.3),
        "Let's build the future of private, intelligent engineering.",
        font_size=12,
        color=ACCENT_BLUE,
        alignment=PP_ALIGN.CENTER,
    )

    add_slide_number(slide, 10, total_slides)

    # ── Save ──────────────────────────────────────────────────────────────
    output_path = Path("/home/dtfrost/neuralmind/docs/market-research/pitch-deck.pptx")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    path = build_deck()
    print(f"✓ Pitch deck saved to: {path}")
    print("  Slides: 10")
    print("  Theme: Dark professional (navy/dark gray, white text, accent blue #0066CC)")
